import csv, json, shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "outputs" / "figures"
RES = ROOT / "outputs" / "results"

results = []
for f in sorted(RES.glob("E*_result.json")):
    results.append(json.loads(f.read_text(encoding="utf-8"))["metrics"])

doc_path = ROOT / "DoAn_Drowsiness_Detection.docx"
ppt_path = ROOT / "DoAn_Slides_BaoVe.pptx"
if not (ROOT / "DoAn_Drowsiness_Detection_original.docx").exists():
    shutil.copy2(doc_path, ROOT / "DoAn_Drowsiness_Detection_original.docx")
if not (ROOT / "DoAn_Slides_BaoVe_original.pptx").exists():
    shutil.copy2(ppt_path, ROOT / "DoAn_Slides_BaoVe_original.pptx")

# Word: replace every explicit image marker with the generated figure.
doc = Document(doc_path)
for p in doc.paragraphs:
    if "[CHÈN HÌNH:" not in p.text:
        continue
    names = [x.name for x in FIG.glob("*.png") if x.name in p.text]
    p.clear()
    if not names:
        p.add_run("[Không có ảnh webcam trong lần chạy tự động local]")
        continue
    for name in names:
        p.add_run().add_picture(str(FIG / name), width=Inches(6.2))
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

hp = doc.add_paragraph()
hr = hp.add_run("PHỤ LỤC A. KẾT QUẢ CHẠY NOTEBOOK (PILOT/LOCAL CPU)")
hr.bold = True; hr.font.size = __import__('docx').shared.Pt(16)
p = doc.add_paragraph(
    "Notebook được chạy ngày 15/08/2026 trên CPU local ở chế độ pilot: lấy mẫu 2% "
    "trong từng tập, ảnh 96×96, 2 epoch phần head và 1 epoch fine-tune. "
    "Các số liệu dưới đây là kết quả chạy thực tế để kiểm tra pipeline, không thay thế "
    "kết quả huấn luyện đầy đủ 224×224 trên Colab GPU."
)
headers = ["TN", "Mô hình", "Chia", "Accuracy", "Precision", "Recall", "F1", "AUC", "Ảnh test"]
t = doc.add_table(rows=1, cols=len(headers))
for i, h in enumerate(headers): t.rows[0].cells[i].text = h
for r in results:
    cells = t.add_row().cells
    vals = [r["Thí nghiệm"], r["Mô hình"], r["Giao thức"], f'{r["Accuracy"]:.4f}',
            f'{r["Precision"]:.4f}', f'{r["Recall"]:.4f}', f'{r["F1-Score"]:.4f}',
            f'{r["ROC-AUC"]:.4f}', str(r["Số ảnh test"])]
    for i, v in enumerate(vals): cells[i].text = v

for title, name in [
    ("So sánh tổng hợp hai giao thức", "hinh05_so_sanh_giao_thuc.png"),
    ("Grad-CAM E2", "hinh06_gradcam_E2.png"),
    ("Grad-CAM E5", "hinh06_gradcam_E5.png"),
]:
    hp = doc.add_paragraph(); rr = hp.add_run(title); rr.bold = True
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(FIG / name), width=Inches(6.2))
doc.save(doc_path)

# PowerPoint: fill key result slides and append a compact run-information slide.
prs = Presentation(ppt_path)

def clear_marker(slide, marker):
    for sh in slide.shapes:
        if hasattr(sh, "text") and marker in sh.text:
            sh.text = ""

def add_pic(slide, name, left, top, width):
    slide.shapes.add_picture(str(FIG / name), PInches(left), PInches(top), width=PInches(width))

# Existing slides are 1-indexed in the presentation narrative.
s12 = prs.slides[11]; clear_marker(s12, "CHÈN:"); add_pic(s12, "hinh_E5_duong_cong_huan_luyen.png", 0.7, 1.6, 8.0)
s14 = prs.slides[13]; clear_marker(s14, "CHÈN:"); add_pic(s14, "hinh05_so_sanh_giao_thuc.png", 0.45, 1.55, 8.3)
s15 = prs.slides[14]; clear_marker(s15, "CHÈN:"); add_pic(s15, "hinh_E2_confusion_matrix.png", 0.35, 1.65, 4.25); add_pic(s15, "hinh_E5_confusion_matrix.png", 4.65, 1.65, 4.25)
s16 = prs.slides[15]; clear_marker(s16, "CHÈN:"); add_pic(s16, "hinh06_gradcam_E5.png", 0.35, 1.55, 4.25); add_pic(s16, "hinh06_gradcam_E2.png", 4.65, 1.55, 4.25)

# Add results table to slide 13.
s13 = prs.slides[12]
rows, cols = 7, 7
shape = s13.shapes.add_table(rows, cols, PInches(0.45), PInches(1.55), PInches(8.6), PInches(4.5))
table = shape.table
hh = ["TN", "Mô hình", "Chia", "Acc", "Prec", "Recall", "F1"]
for j, h in enumerate(hh): table.cell(0,j).text = h
for i, r in enumerate(results, 1):
    vals = [r["Thí nghiệm"], r["Mô hình"], r["Giao thức"], f'{r["Accuracy"]:.3f}', f'{r["Precision"]:.3f}', f'{r["Recall"]:.3f}', f'{r["F1-Score"]:.3f}']
    for j, v in enumerate(vals): table.cell(i,j).text = v
for row in table.rows:
    for cell in row.cells:
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(10); para.alignment = PP_ALIGN.CENTER

layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
title_box = slide.shapes.add_textbox(PInches(0.8), PInches(0.45), PInches(8.0), PInches(0.7))
title_box.text_frame.text = "Thông tin lần chạy notebook"
title_box.text_frame.paragraphs[0].font.size = Pt(30)
title_box.text_frame.paragraphs[0].font.bold = True
box = slide.shapes.add_textbox(PInches(0.8), PInches(1.6), PInches(8.0), PInches(4.7))
tf = box.text_frame
lines = [
    "• Thời điểm: 15/08/2026; môi trường local CPU; TensorFlow 2.21",
    "• Chế độ pilot: 2% mỗi tập, ảnh 96×96, 2 + 1 epoch",
    "• Đã chạy đủ E1–E6 và xuất đường cong, confusion matrix, ROC, Grad-CAM",
    "• Kết quả pilot chỉ dùng kiểm tra pipeline; cần chạy lại cấu hình đầy đủ trên Colab GPU trước khi công bố số liệu chính thức.",
]
tf.text = lines[0]
for line in lines[1:]: tf.add_paragraph().text = line
for p in tf.paragraphs: p.font.size = Pt(22); p.space_after = Pt(14)
prs.save(ppt_path)

print(f"Updated {doc_path.name} and {ppt_path.name}")
